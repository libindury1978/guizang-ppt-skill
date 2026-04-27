#!/usr/bin/env python3
"""Build editable PPTX from a JSON deck spec.

Usage:
  python3 scripts/build_pptx.py --spec path/to/deck.json --output path/to/deck.pptx
"""

from __future__ import annotations

import argparse
import json
from pathlib import Path
from typing import Any, Dict, List

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


THEMES = {
    "ink-classic": {"title": RGBColor(33, 33, 33), "body": RGBColor(70, 70, 70)},
    "indigo-porcelain": {"title": RGBColor(40, 60, 110), "body": RGBColor(65, 85, 130)},
    "forest-ink": {"title": RGBColor(40, 80, 60), "body": RGBColor(60, 105, 85)},
    "kraft-paper": {"title": RGBColor(85, 60, 35), "body": RGBColor(115, 85, 60)},
    "dune": {"title": RGBColor(65, 55, 45), "body": RGBColor(95, 80, 65)},
}


def add_textbox(slide, text: str, left: float, top: float, width: float, height: float, size: int, bold=False, color=None, align=PP_ALIGN.LEFT):
    shape = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text or ""
    run.font.size = Pt(size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = color
    p.alignment = align
    return shape


def add_cover(slide, item: Dict[str, Any], colors):
    add_textbox(slide, item.get("title", ""), 0.8, 1.7, 11.0, 1.2, 42, bold=True, color=colors["title"])
    add_textbox(slide, item.get("subtitle", ""), 0.8, 3.2, 10.5, 1.0, 22, color=colors["body"])
    add_textbox(slide, item.get("author", ""), 0.8, 6.5, 6.0, 0.5, 14, color=colors["body"])


def add_section(slide, item: Dict[str, Any], colors):
    add_textbox(slide, item.get("title", ""), 1.0, 2.3, 11.0, 1.0, 36, bold=True, color=colors["title"])
    add_textbox(slide, item.get("subtitle", ""), 1.0, 3.5, 10.0, 0.8, 20, color=colors["body"])


def add_content_image(slide, item: Dict[str, Any], colors, base_dir: Path):
    add_textbox(slide, item.get("title", ""), 0.7, 0.6, 5.8, 0.8, 30, bold=True, color=colors["title"])
    bullets: List[str] = item.get("bullets", [])
    top = 1.6
    for b in bullets[:6]:
        add_textbox(slide, f"• {b}", 0.9, top, 5.2, 0.5, 18, color=colors["body"])
        top += 0.58
    image = item.get("image")
    if image:
        image_path = (base_dir / image).resolve()
        if image_path.exists():
            slide.shapes.add_picture(str(image_path), Inches(6.3), Inches(1.2), width=Inches(6.6))
        else:
            add_textbox(slide, f"[missing image] {image}", 6.4, 3.2, 5.8, 0.7, 14, color=colors["body"], align=PP_ALIGN.CENTER)


def add_quote(slide, item: Dict[str, Any], colors):
    quote = item.get("quote") or item.get("title", "")
    add_textbox(slide, f"“{quote}”", 1.0, 2.2, 11.3, 2.0, 34, bold=True, color=colors["title"], align=PP_ALIGN.CENTER)
    add_textbox(slide, item.get("source", ""), 1.0, 5.0, 11.3, 0.6, 16, color=colors["body"], align=PP_ALIGN.CENTER)


def add_comparison(slide, item: Dict[str, Any], colors):
    add_textbox(slide, item.get("title", "Before / After"), 0.7, 0.6, 11.8, 0.7, 30, bold=True, color=colors["title"])
    left = item.get("left", [])
    right = item.get("right", [])
    add_textbox(slide, "Before", 0.8, 1.5, 5.4, 0.5, 20, bold=True, color=colors["body"])
    add_textbox(slide, "After", 6.9, 1.5, 5.4, 0.5, 20, bold=True, color=colors["body"])
    top = 2.1
    for t in left[:6]:
        add_textbox(slide, f"• {t}", 0.9, top, 5.2, 0.45, 16, color=colors["body"])
        top += 0.52
    top = 2.1
    for t in right[:6]:
        add_textbox(slide, f"• {t}", 7.0, top, 5.2, 0.45, 16, color=colors["body"])
        top += 0.52


def add_stats(slide, item: Dict[str, Any], colors):
    add_textbox(slide, item.get("title", "关键数据"), 0.8, 0.6, 11.4, 0.7, 30, bold=True, color=colors["title"])
    stats = item.get("stats", [])
    x = 0.9
    for stat in stats[:3]:
        add_textbox(slide, str(stat.get("value", "")), x, 2.1, 3.6, 1.0, 36, bold=True, color=colors["title"], align=PP_ALIGN.CENTER)
        add_textbox(slide, stat.get("label", ""), x, 3.2, 3.6, 0.6, 16, color=colors["body"], align=PP_ALIGN.CENTER)
        x += 4.1


def add_pipeline(slide, item: Dict[str, Any], colors):
    add_textbox(slide, item.get("title", "流程"), 0.8, 0.6, 11.4, 0.7, 30, bold=True, color=colors["title"])
    steps = item.get("steps", [])
    y = 1.8
    i = 1
    for step in steps[:6]:
        add_textbox(slide, f"{i}. {step}", 1.0, y, 10.8, 0.5, 18, color=colors["body"])
        y += 0.72
        i += 1


def add_image_grid(slide, item: Dict[str, Any], colors, base_dir: Path):
    add_textbox(slide, item.get("title", "图集"), 0.8, 0.6, 11.4, 0.7, 30, bold=True, color=colors["title"])
    images = item.get("images", [])[:4]
    positions = [(0.9, 1.6), (6.7, 1.6), (0.9, 4.0), (6.7, 4.0)]
    for img, pos in zip(images, positions):
        p = (base_dir / img).resolve()
        if p.exists():
            slide.shapes.add_picture(str(p), Inches(pos[0]), Inches(pos[1]), width=Inches(5.0), height=Inches(2.1))
        else:
            add_textbox(slide, f"[missing] {img}", pos[0], pos[1] + 0.8, 5.0, 0.5, 12, color=colors["body"], align=PP_ALIGN.CENTER)


def build(spec: Dict[str, Any], output: Path, base_dir: Path):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    theme_name = spec.get("theme", "ink-classic")
    colors = THEMES.get(theme_name, THEMES["ink-classic"])

    for item in spec.get("slides", []):
        slide = prs.slides.add_slide(blank)
        typ = item.get("type", "content_image")
        if typ == "cover":
            add_cover(slide, item, colors)
        elif typ == "section":
            add_section(slide, item, colors)
        elif typ == "quote":
            add_quote(slide, item, colors)
        elif typ == "comparison":
            add_comparison(slide, item, colors)
        elif typ == "stats":
            add_stats(slide, item, colors)
        elif typ == "pipeline":
            add_pipeline(slide, item, colors)
        elif typ == "image_grid":
            add_image_grid(slide, item, colors, base_dir)
        else:
            add_content_image(slide, item, colors, base_dir)

    output.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output))


def parse_args():
    parser = argparse.ArgumentParser(description="Generate editable PPTX from deck JSON spec.")
    parser.add_argument("--spec", required=True, help="Path to deck json spec.")
    parser.add_argument("--output", required=True, help="Path to output pptx.")
    return parser.parse_args()


def main():
    args = parse_args()
    spec_path = Path(args.spec).expanduser().resolve()
    output = Path(args.output).expanduser().resolve()
    with spec_path.open("r", encoding="utf-8") as f:
        spec = json.load(f)
    build(spec, output, spec_path.parent)
    print(f"Generated: {output}")


if __name__ == "__main__":
    main()
